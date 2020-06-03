# test_extract.py

import subprocess # pip install --pre pyhwp
from pytesseract import *  # pip install pytesseract
import cv2  # pip install opencv-python
import docx2txt  # pip install docx2txt
import pdfplumber  # pip install pdfplumber
from pptx import Presentation  # pip install python-pptx

# 확장자별로 이미지를 불러오는 방식을 다르게 함

# 모든 이미지 파일에 사용되는 함수
def ImagetoText(fileName):
    # 이미지 불러오기
    image = cv2.imread(fileName)
    # 이미지를 흑백으로 변경
    image = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    # 이미지를 흐릿하게 (가우시안 블러)
    image = cv2.GaussianBlur(image, (3, 3), 0)
    # 이미지의 노이즈 줄이기
    image = cv2.fastNlMeansDenoising(image, h=10, searchWindowSize=21, templateWindowSize=7)
    # 이미지를 흰색과 검은색으로 임계 전처리
    image_final = cv2.threshold(image, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
    # tesseract 옵션 설정, 언어:한글 (+세로쓰기 글자에는 Hangul_vert), psm = 3, 띄어쓰기 보완

    config = "-l script/Hangul --oem 1 --psm 3 -c preserve_interword_spaces=1"

    # 텍스트 추출
    extracted_text = image_to_string(image_final, config=config)
    if extracted_text == "":  # 텍스트가 없는 이미지이거나 인식이 안 되는 이미지의 경우
        extracted_text = "텍스트를 발견하지 못했습니다."
    return extracted_text

# txt 파일에만 사용되는 함수
def TxttoText(fileName):
    text = open(fileName, mode="rt", encoding="utf-8")
    extracted_text = text.read()
    if extracted_text == "":  # 빈 텍스트 파일일 경우
        extracted_text = "텍스트를 발견하지 못했습니다."
    return extracted_text

# doc, docx 파일에만 사용되는 함수
def DocxtoText(fileName):
    text = docx2txt.process(fileName)
    extracted_text = text
    if extracted_text == "":  # 빈 텍스트 파일일 경우
        extracted_text = "텍스트를 발견하지 못했습니다."
    return extracted_text

# pdf 파일에만 사용되는 함수
def PdftoText(fileName):
    pdf = pdfplumber.open(fileName)
    extracted_text = ""
    for i in range(0, len(pdf.pages)):
        text = pdf.pages[i].extract_text()
        if type(text) == str:
            extracted_text += "".join(text)+"\n"
    if extracted_text == "":  # 빈 텍스트 파일일 경우
        extracted_text = "텍스트를 발견하지 못했습니다."
    return extracted_text

# ppt, pptx 파일에만 사용되는 함수
def PptxtoText(fileName):
    text = Presentation(fileName)
    extracted_text = ""
    for slide in text.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                extracted_text += "".join(paragraph.text)+"\n"
    return extracted_text

# hwp 파일에만 사용되는 함수
def HwptoText(fileName):
    text = subprocess.check_output(['hwp5txt', fileName])
    extracted_text = text.decode("utf-8")
    return extracted_text

def ReturnText(fileName):  # switch 문이 없어서 우선 if-else로 작성
    filetype = fileName.split(".")[-1]
    filetype = filetype.lower()
    if filetype in ["bmp", "jpg", "jpeg", "png"]:  # 이미지 파일인 경우 (필요에 따라 확장자 추가)
        text = ImagetoText(fileName)
    elif filetype == "txt":  # .txt 파일인 경우
        text = TxttoText(fileName)
    elif filetype == "docx":  # .docx 파일인 경우 (.doc은 x)
        text = DocxtoText(fileName)
    elif filetype == "pdf":   # .pdf 파일인 경우
        text = PdftoText(fileName)
    elif filetype == "pptx":  # .pptx 파일인 경우 (.ppt은 x)
        text = PptxtoText(fileName)
    elif filetype == "hwp":  # .hwp 파일인 경우
        text = HwptoText(fileName)
    else:  # mp3, zip 등 지원하지 않는 확장자인 경우
        text = "변환할 수 없는 파일입니다.\n지원하는 파일 타입은 이미지 파일 또는 텍스트 파일입니다.\n다시 시도해 주십시오."
    return text

if __name__ == "__main__":
    print("this is text_extract.py")